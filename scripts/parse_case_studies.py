#!/usr/bin/env python3
"""
Parse assets/ case studies (PPTX), extract structured content.
Output: Case study → Service, Pain, Outcome, and mapping to cadences.
"""

import json
from pathlib import Path
from pptx import Presentation

BASE = Path(__file__).resolve().parent.parent
ASSETS = BASE / "assets"
OUT = BASE / "reports" / "Case Study Cadence Mapping.md"


def extract_pptx_text(path: Path) -> list[dict]:
    """Extract slide title + body text from PPTX."""
    slides = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if texts:
            slides.append({"slide": i + 1, "text": " | ".join(texts)})
    return slides


def parse_case_studies() -> list[dict]:
    """Parse all PPTX in assets and return structured records."""
    records = []
    search = ASSETS / "case-studies" if (ASSETS / "case-studies").exists() else ASSETS
    files = sorted(search.rglob("*.pptx")) if search.is_dir() else []
    for f in files[:50]:  # limit for cadence mapping; full set in parse_all_case_studies
        slides = extract_pptx_text(f)
        full_text = " ".join(s["text"] for s in slides).lower()
        # Infer client, service, pain patterns
        client = f.stem.split(" – ")[0].split(" - ")[0].replace(".pptx", "").strip()
        if "microsoft" in client.lower() or "msft" in client.lower():
            client = "Microsoft"
        elif "saint" in client.lower() or "st. luke" in client.lower():
            client = "St. Luke's"
        elif "alorica" in client.lower():
            client = "Alorica"

        service = "AI / Automation"  # default for most
        if "zerotrust" in f.name.lower() or "webinar" in full_text:
            service = "Creative"
        elif "jfk" in f.name.lower() or "search" in full_text and "cognitive" in full_text:
            service = "Data & Analytics"

        pain_keywords = []
        if "manual" in full_text or "manually" in full_text:
            pain_keywords.append("manual processes")
        if "repetitive" in full_text or "repetitive" in full_text:
            pain_keywords.append("repetitive inquiries")
        if "scattered" in full_text or "multiple systems" in full_text or "multiple data" in full_text:
            pain_keywords.append("scattered knowledge/data")
        if "time" in full_text and ("intensive" in full_text or "consuming" in full_text):
            pain_keywords.append("time-intensive")
        if "duplicate" in full_text or "re-entered" in full_text:
            pain_keywords.append("duplicate data entry")
        if "excel" in full_text:
            pain_keywords.append("Excel/spreadsheet dependence")
        if "convert" in full_text and "lead" in full_text:
            pain_keywords.append("lead conversion")
        if "rfp" in full_text:
            pain_keywords.append("RFP response speed")
        if "customer service" in full_text or "support" in full_text:
            pain_keywords.append("support volume/resolution")
        if "hr" in full_text or "benefits" in full_text or "pto" in full_text:
            pain_keywords.append("HR inquiry volume")
        if "compliance" in full_text or "regulated" in full_text:
            pain_keywords.append("compliance/regulated workflows")
        if "event" in full_text and "calendar" in full_text:
            pain_keywords.append("event data collection/prioritization")

        records.append({
            "file": f.name,
            "client": client,
            "service": service,
            "pain_keywords": list(set(pain_keywords)) if pain_keywords else ["automation"],
            "slides": len(slides),
            "summary": slides[0]["text"][:400] if slides else "",
        })
    return records


def main():
    records = parse_case_studies()
    # Write markdown report
    lines = [
        "# Case Study → Cadence Mapping",
        "",
        "**Source:** assets/ (8 case study decks)  ",
        "**Purpose:** Map case studies to outbound cadences and pains for targeted proof.",
        "",
        "---",
        "",
        "## 1. Parsed Case Studies",
        "",
        "| Client | Service | Pains Addressed | File |",
        "|--------|---------|-----------------|------|",
    ]
    for r in records:
        pains = ", ".join(r["pain_keywords"][:5])
        lines.append(f"| {r['client']} | {r['service']} | {pains} | `{r['file']}` |")

    lines.extend([
        "",
        "---",
        "",
        "## 2. Case Study → Cadence Campaign Mapping",
        "",
        "Use the case study that matches the **pain** your prospect has. Attach or link in cadence emails where CTA or proof is relevant.",
        "",
        "| Case Study | Best Cadence(s) | Use When Prospect Says... | Cadence Days |",
        "|------------|-----------------|---------------------------|--------------|",
    ])

    mapping = [
        ("Alorica (Website Chat + RFP)", "01 Marketing Creative, 02 Attribution", "conversion, lead quality, RFP response, sales enablement", "5, 10, 15, 20"),
        ("CoPilot HR Agent", "07 Finance Director, 05 Decision Latency", "HR questions, benefits, PTO, scattered policies, manual responses", "8–14"),
        ("CoPilot Customer Service Agent", "03 IT Director Ownership, 04 Drift", "support volume, knowledge scattered, slow resolution, repetitive tickets", "5, 10, 15"),
        ("CoPilot Financial Analysis Agent", "05 Decision Latency, 10 C-level KPI", "time-to-answer, manual reporting, decision speed", "1, 5, 15"),
        ("Microsoft Partner Support Agent", "08 IT Director Backlog, 03 Ownership", "manual support, multiple data sources, partner/incentive questions", "11, 15"),
        ("Microsoft Pharmaceutical AI PoC", "05 Decision Latency, 10 C-level", "compliance, regulated workflows, R&D, time-to-insight", "1–7"),
        ("Microsoft JFK Files / AI Search", "05 Decision Latency, Data & Analytics", "search, findability, data discovery", "8–14"),
        ("Microsoft Bing Calendar Agent", "05 Decision Latency, 08 Backlog", "manual data collection, Excel, event prioritization, scalability", "11, 20"),
        ("St. Luke's Medical Scribing", "05 Decision Latency, 08 Backlog", "duplicate entry, manual EMR, clinical documentation, healthcare efficiency", "8–14"),
        ("MSFT ZeroTrust Webinar", "01 Marketing Creative (video/delivery)", "webinar, demos, delivery, testimonial", "N/A (Creative asset)"),
    ]

    for casestudy, cadence, trigger, days in mapping:
        lines.append(f"| {casestudy} | {cadence} | {trigger} | {days} |")

    lines.extend([
        "",
        "---",
        "",
        "## 3. Pain → Service → Case Study",
        "",
        "| Pain | Affirma Service | Case Study to Use |",
        "|------|-----------------|-------------------|",
    ])

    pain_map = [
        ("Manual processes / duplicate data entry", "AI / Automation", "St. Luke's, Microsoft Partner Support, Bing Calendar"),
        ("Repetitive inquiries (HR, Support, Partners)", "AI / Automation", "CoPilot HR, Customer Service, Partner Support"),
        ("Scattered knowledge / multiple systems", "AI / Automation, Modern Workplace", "CoPilot agents, Alorica RFP"),
        ("Time-to-answer / decision latency", "Data & Analytics, AI / Automation", "Financial Analysis, Pharmaceutical, JFK Search"),
        ("Lead conversion / website discovery", "Marketing Services, Creative", "Alorica Website Chat"),
        ("RFP response speed / sales enablement", "AI / Automation, Custom Dev", "Alorica RFP Assistant"),
        ("Event/data collection, Excel dependence", "AI / Automation", "Microsoft Bing Calendar"),
        ("Healthcare/clinical documentation", "AI / Automation", "St. Luke's Medical Scribing"),
        ("Compliance / regulated workflows", "AI / Automation", "Microsoft Pharmaceutical"),
    ]

    for pain, svc, cs in pain_map:
        lines.append(f"| {pain} | {svc} | {cs} |")

    lines.extend([
        "",
        "---",
        "",
        "## 4. Cadence Day Recommendations",
        "",
        "**When to attach a case study:**",
        "- **Days 5, 10, 15, 20, 25, 30** — Soft CTA days. Attach 1 relevant case study with the doc/checklist.",
        "- **Days 8–14** — Deeper value phase. Use case study as 'proof' in body: 'We did X for [Client] — similar situation.'",
        "- **Day 1** — Only if prospect already signaled the pain. Don't lead with case study; establish pattern first.",
        "",
        "**Rule:** Match case study pain to cadence theme. IT Director Drift → use Customer Service or Partner Support (scattered knowledge). Marketing Creative → use Alorica (conversion).",
        "",
        "---",
        "",
        "**Regenerate:** `python3 scripts/parse_case_studies.py`",
    ])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {OUT}")
    print(f"Parsed {len(records)} case studies")


if __name__ == "__main__":
    main()
