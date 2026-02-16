#!/usr/bin/env python3
"""
Deep parse of commit 8063f246 files.
"""

import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pypdf import PdfReader
import json

OUT = []

def log(s):
    OUT.append(s)
    print(s)

def parse_xlsx(path, max_rows=20, max_cols=30):
    """Parse Excel, show structure and sample."""
    try:
        xl = pd.ExcelFile(path)
        for sheet in xl.sheet_names:
            df = pd.read_excel(path, sheet_name=sheet, header=None)
            log(f"\n{'='*60}")
            log(f"FILE: {path.name} | Sheet: {sheet}")
            log(f"Shape: {df.shape[0]} rows × {df.shape[1]} cols")
            if df.shape[0] > 0:
                # First row might be header
                log("First 3 rows (raw):")
                for i in range(min(3, len(df))):
                    log(str(df.iloc[i].tolist()[:max_cols]))
                if len(df) > 3:
                    log("...")
        return True
    except Exception as e:
        log(f"ERROR {path}: {e}")
        return False

def parse_pptx(path):
    """Extract text from all slides."""
    try:
        prs = Presentation(path)
        log(f"\n{'='*60}")
        log(f"FILE: {path.name} | Slides: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = shape.text.strip()
                    if t:
                        texts.append(t)
            if texts:
                log(f"  Slide {i+1}:")
                for t in texts[:5]:  # first 5 text blocks
                    log(f"    {t[:200]}{'...' if len(t)>200 else ''}")
        return True
    except Exception as e:
        log(f"ERROR {path}: {e}")
        return False

def parse_pdf(path):
    """Extract text from PDF."""
    try:
        reader = PdfReader(path)
        log(f"\n{'='*60}")
        log(f"FILE: {path.name} | Pages: {len(reader.pages)}")
        for i, page in enumerate(reader.pages[:5]):  # first 5 pages
            t = page.extract_text()
            if t:
                lines = t.strip().split('\n')[:15]
                log(f"  Page {i+1} (first 15 lines):")
                for line in lines:
                    log(f"    {line[:100]}")
        return True
    except Exception as e:
        log(f"ERROR {path}: {e}")
        return False

def main():
    # Run from repo root
    base = Path(__file__).resolve().parent.parent
    data_dir = base / "data"
    assets_dir = base / "assets"
    
    # Excel files (data/opportunities and data/accounts)
    log("=" * 80)
    log("DEEP PARSE — Commit 8063f246 (chehaba/prospecting)")
    log("=" * 80)
    
    opp_dir = data_dir / "opportunities"
    acct_dir = data_dir / "accounts"
    for name, subdir in [
        ("All Account 2-12-2026 1-27-59 PM.xlsx", acct_dir),
        ("All Opportunities 2-11-2026 12-50-11 PM.xlsx", opp_dir),
        ("Pipeline Management(All Opportunities Export).csv.xlsx", opp_dir),
    ]:
        p = subdir / name
        if p.exists():
            parse_xlsx(p)
    
    # Numbers - Apple format, try unzip
    num_path = opp_dir / "Consolidated_Opportunities 2.numbers"
    if num_path.exists():
        log(f"\n{'='*60}")
        log(f"FILE: {num_path.name}")
        log("Numbers is Apple format; structure requires Numbers or conversion. Skipping raw parse.")
    
    # PowerPoint files (key case studies; search in case-studies or assets)
    pptx_dir = assets_dir / "case-studies" if (assets_dir / "case-studies").exists() else assets_dir
    key_names = [
        "Alorica – Website Chat Bot & RFP Assistant.pptx",
        "CoPilot Case Studies.pptx",
        "MSFT_ZeroTrust-Onepage-CS.pptx",
        "Microsoft - Partner Support AI Agent .pptx",
        "Saint Lukes - Medical Scribing AI Agent.pptx",
    ]
    for p in pptx_dir.rglob("*.pptx"):
        if p.name in key_names:
            parse_pptx(p)
    
    # PDF
    ref_dir = data_dir / "reference"
    pdf_path = ref_dir / "System%20Prompt.pdf.pdf"
    if pdf_path.exists():
        parse_pdf(pdf_path)
    
    # Consolidated_Opportunities.pdf (we've parsed before)
    pdf2 = opp_dir / "Consolidated_Opportunities.pdf"
    if pdf2.exists():
        log(f"\n{'='*60}")
        log(f"FILE: Consolidated_Opportunities.pdf (already parsed in prior analysis)")
    
    return "\n".join(OUT)

if __name__ == "__main__":
    main()
