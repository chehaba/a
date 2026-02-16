#!/usr/bin/env python3
"""
Parse all assets, extract client + service from filenames and (where possible) content.
Output: Value prop mapping — "We help [X] clients like [Client] do [Y] which resulted in [Z]"
"""

import re
from pathlib import Path
from collections import defaultdict

BASE = Path(__file__).resolve().parent.parent
ASSETS = BASE / "assets"
OUT = BASE / "reports" / "Case Study Value Prop Mapping.md"

# Filename patterns: Client - Service or Client – Service (em dash)
SERVICE_KEYWORDS = {
    "Modern Workplace": ["sharepoint", "spo", "m365", "o365", "teams", "intranet", "migration", "onedrive", "exchange"],
    "CRM": ["crm", "salesforce", "dynamics", "d365", "bc ", " business central", "cpq", "sales "],
    "Data & Analytics": ["bi", "power bi", "snowflake", "data warehouse", "analytics", "d&a", "dashboard"],
    "Marketing": ["marketing", "hubspot", "pardot", "marketo", "ppc", "seo", "sem", "demand gen", "social media", "abm"],
    "Custom Dev": ["custom dev", "mobile", "app", "portal", "application", "development", "integration"],
    "AI / Automation": ["ai ", "copilot", "agent", "automation", "machine learning"],
    "Infrastructure": ["infrastructure", "infra", "mdm", "intune", "sccm", "exchange", "bp [o]"],
    "Creative": ["creative", "video", "webinar", "graphic", "design", "recording", "signage", "catalog"],
}

INDUSTRY_FROM_CLIENT = {
    r"st\.?\s*luke|saint lukes|hospital|medical|health|evergreen|kaiser|childrens|marin general|burkhart dental|casey family": "Healthcare",
    r"city of |county of |arapahoe|ramsey county|port of tacoma|tualatin valley|scottsdale|hillsboro|beaverton|medina|seattle|state parks|weci": "Government",
    r"university|college|colorado mountain|oregon state|ucla|nii mbl|lake county|thales academy": "Education",
    r"credit union|bank|insurance|selco|peoples bank|solarity|chevron credit|hawaii bank|spectrum credit|alaska national": "Financial Services",
    r"foundation|cystic fibrosis|mary's place|degrees of change|ymca|electronic frontier|art museum": "Non-Profit",
    r"ferguson|dupont|snap on|resco|avista|hermanson|samson rope|hydac|hasco|captiveaire": "Manufacturing / Industrial",
    r"agwest|nw farm credit|farm credit": "Financial Services / Agriculture",
}


def extract_client_service(path: Path) -> tuple[str, str, list[str]]:
    """From path like 'St. Lukes - Medical Scribing AI Agent.pptx' -> (client, primary_service, all_services)"""
    name = path.stem.replace("–", "-").replace("—", "-")
    # Split on " - " or " – "
    parts = re.split(r"\s+[-–—]\s+", name, maxsplit=1)
    client = parts[0].strip() if parts else name
    rest = parts[1].strip() if len(parts) > 1 else ""

    # Clean client (remove "Case Study", "Approved", etc.)
    client = re.sub(r"\s*-\s*Approved.*$", "", client, flags=re.I)
    client = re.sub(r"\s*Case Study.*$", "", client, flags=re.I)
    client = re.sub(r"\s*\(.*\)$", "", client)
    client = client.strip()

    detected_services = []
    lower_rest = rest.lower() + " " + path.name.lower()
    # Strong filename cues (check first - more specific)
    if "sharepoint" in lower_rest or " spo " in lower_rest or " sp " in lower_rest or " spo migration" in lower_rest or "m365" in lower_rest or "o365" in lower_rest or "teams" in lower_rest or "intranet" in lower_rest:
        detected_services.append("Modern Workplace")
    if "salesforce" in lower_rest or "dynamics" in lower_rest or " d365" in lower_rest or " bc " in lower_rest or "business central" in lower_rest or " crm " in lower_rest:
        detected_services.append("CRM")
    if " bi " in lower_rest or " - bi" in lower_rest or "power bi" in lower_rest or "snowflake" in lower_rest or "data warehouse" in lower_rest or "dashboard" in lower_rest or "d&a" in lower_rest:
        detected_services.append("Data & Analytics")
    if "marketing" in lower_rest or "hubspot" in lower_rest or "pardot" in lower_rest or "marketo" in lower_rest or "ppc" in lower_rest or "seo" in lower_rest or "demand gen" in lower_rest:
        detected_services.append("Marketing")
    if "custom dev" in lower_rest or "mobile" in lower_rest or " application" in lower_rest or "portal" in lower_rest or "development" in lower_rest:
        detected_services.append("Custom Dev")
    if " ai " in lower_rest or "copilot" in lower_rest or "agent" in lower_rest or "automation" in lower_rest:
        detected_services.append("AI / Automation")
    if "infrastructure" in lower_rest or "mdm" in lower_rest or "intune" in lower_rest or "sccm" in lower_rest or "exchange" in lower_rest:
        detected_services.append("Infrastructure")
    if "creative" in lower_rest or "video" in lower_rest or "webinar" in lower_rest or "graphic" in lower_rest or "recording" in lower_rest or "signage" in lower_rest:
        detected_services.append("Creative")
    if "migration" in lower_rest and "Modern Workplace" not in detected_services and not any(s in detected_services for s in ["CRM", "Marketing"]):
        detected_services.append("Modern Workplace")
    # Fallback for Other: weaker signals
    if not detected_services:
        if "website" in lower_rest or "sitefinity" in lower_rest or "site " in lower_rest:
            detected_services.append("Creative")
        elif "exchange" in lower_rest or "dlp" in lower_rest:
            detected_services.append("Infrastructure")
        elif "sharepoint" in path.name.lower():  # sometimes in client or path
            detected_services.append("Modern Workplace")
        elif "governance" in lower_rest:
            detected_services.append("Modern Workplace")
        elif "contract" in lower_rest or "proposal" in lower_rest:
            detected_services.append("Other")
    primary = detected_services[0] if detected_services else "Other"
    return client, primary, detected_services


def infer_industry(client: str) -> str:
    c = client.lower()
    for pat, ind in INDUSTRY_FROM_CLIENT.items():
        if re.search(pat, c):
            return ind
    return "General"


def extract_pptx_snippet(path: Path, max_chars: int = 300) -> str:
    """Extract first slide text from PPTX for outcome hint."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        if not prs.slides:
            return ""
        texts = []
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        full = " ".join(texts)
        # Look for Results, Outcome, Business Need
        for marker in ["Results", "Outcome", "What We Did", "resulted", "reduced", "increased", "%"]:
            if marker.lower() in full.lower():
                break
        return full[:max_chars].replace("\n", " ") if full else ""
    except Exception:
        return ""


def main():
    files = sorted((ASSETS / "case-studies").rglob("*")) if (ASSETS / "case-studies").exists() else sorted(ASSETS.rglob("*"))
    files = [f for f in files if f.suffix.lower() in (".pptx", ".docx", ".pdf") and f.is_file()]

    records = []
    for path in files:
        rel = path.relative_to(ASSETS)
        client, primary, services = extract_client_service(path)
        industry = infer_industry(client)
        snippet = ""
        if path.suffix.lower() == ".pptx":
            snippet = extract_pptx_snippet(path)
        records.append({
            "path": str(rel),
            "client": client,
            "industry": industry,
            "primary": primary,
            "services": services,
            "snippet": snippet,
        })

    # Action verbs by service
    SERVICE_ACTIONS = {
        "Modern Workplace": "migrate SharePoint/M365, build intranets, consolidate tenants",
        "CRM": "implement Salesforce/Dynamics, redesign CRM, integrate sales systems",
        "Data & Analytics": "deploy Power BI, build data warehouses, create dashboards",
        "Marketing": "run PPC/SEO, migrate HubSpot/Marketo, manage demand gen",
        "Custom Dev": "build custom apps, develop portals, create mobile applications",
        "AI / Automation": "deploy AI agents, automate manual processes",
        "Infrastructure": "deploy MDM/Intune, manage Exchange, modernize infrastructure",
        "Creative": "produce video, design assets, deliver webinar content",
        "Other": "deliver technology solutions",
    }

    def value_prop(r):
        industry = r["industry"]
        client = r["client"]
        svc = r["primary"]
        action = SERVICE_ACTIONS.get(svc, "deliver solutions")
        outcome = ""
        if r["snippet"]:
            s = r["snippet"].lower()
            if "reduced" in s or "reduction" in s or "%" in s:
                outcome = "reduced manual effort and improved efficiency"
            elif "migration" in s or "migrat" in s:
                outcome = "successful migration with minimal disruption"
            elif "automation" in s or "agent" in s:
                outcome = "automated manual processes"
            elif "integration" in s:
                outcome = "seamless integration and streamlined workflows"
        if not outcome:
            outcome = "improved operations and measurable outcomes"
        return f"We help {industry} clients like {client} {action} which resulted in {outcome}"

    # Group by service
    by_service = defaultdict(list)
    for r in records:
        by_service[r["primary"]].append(r)

    # Build markdown
    lines = [
        "# Case Study Value Prop Mapping",
        "",
        "**Source:** assets/ ({} files)  ".format(len(records)),
        "**Format:** We help [industry] clients like [Client] do [service] which resulted in [outcome]",
        "",
        "---",
        "",
    ]

    for svc in sorted(by_service.keys()):
        rows = by_service[svc]
        lines.append(f"## {svc} ({len(rows)} case studies)")
        lines.append("")
        lines.append("| Client | Industry | Value Proposition | Asset |")
        lines.append("|--------|----------|-------------------|-------|")
        for r in sorted(rows, key=lambda x: x["client"]):
            vp = value_prop(r)
            # Truncate vp if very long
            vp_short = vp[:100] + "…" if len(vp) > 100 else vp
            link = f"[{r['path']}](../assets/{r['path'].replace(' ', '%20')})"
            lines.append(f"| {r['client']} | {r['industry']} | {vp_short} | {link} |")
        lines.append("")

    # Append full value props section
    lines.extend([
        "---",
        "",
        "## Full Value Propositions (Copy-Ready)",
        "",
    ])
    for svc in sorted(by_service.keys()):
        lines.append(f"### {svc}")
        lines.append("")
        for r in sorted(by_service[svc], key=lambda x: x["client"]):
            vp = value_prop(r)
            lines.append(f"- **{r['client']}:** {vp} — [Case Study](../assets/{r['path'].replace(' ', '%20')})")
        lines.append("")

    # Cadence mapping
    SVC_TO_CADENCE = {
        "Modern Workplace": "03 IT Director Ownership, 04 IT Director Drift",
        "CRM": "Segment Packs (CRM frame); no dedicated cadence",
        "Data & Analytics": "05 Director Decision Latency, 10 C-level KPI",
        "Marketing": "01 Marketing Creative, 02 Attribution, 06 Marketing VP Spend, 09 Marketing Manager",
        "Custom Dev": "08 IT Director Backlog",
        "AI / Automation": "05 Decision Latency, 08 Backlog, 03 Ownership (scattered knowledge)",
        "Infrastructure": "03 IT Director Ownership, 04 Drift",
        "Creative": "01 Marketing Creative",
        "Other": "Match pain to cadence",
    }
    lines.extend([
        "---",
        "",
        "## Service → Outbound Cadence",
        "",
        "| Service | Use in Cadence(s) |",
        "|---------|-------------------|",
    ])
    for svc, cad in SVC_TO_CADENCE.items():
        lines.append(f"| {svc} | {cad} |")
    lines.append("")
    lines.append("**Rule:** When attaching a case study, match the prospect's pain to the case study outcome. Use the value prop in the email body: *\"We helped [similar client] [action] which resulted in [outcome].\"*")
    lines.append("")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"Parsed {len(records)} assets")
    print(f"By service: {dict((k, len(v)) for k, v in sorted(by_service.items()))}")
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
